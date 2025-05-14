import os
import re
from collections import defaultdict

"""
구절 입력 규칙
1) 각 구절은 '숫자. 책 장:절' 형식으로 입력합니다.
2) 한 슬라이드에 여러 구절을 넣고 싶다면, 각 구절을 쉼표로 구분합니다.
3) 이어진 구절로 이루어져 있다면 '책 장:절-절' 형식으로 입력합니다.
4) 책 제목은 약어로 입력합니다. 예) '1. 잠 1:8, 잠 31:2' 또는 '2. 잠 4:1-9'

예)
1. 잠 1:8, 잠 31:2
2. 마 5:7-9
"""

raw_text = """
1. 잠 1:8, 잠 31:2
2. 잠 4:1, 5:7
3. 잠 4:1-9
4. 잠 22:6, 잠 22:15
5. 잠 1:7-9
6. 잠 4:7-9
7. 잠 22:3
8. 잠 22:5-6
9. 잠 22:4
10. 잠 3:34-35
11. 잠 22:4
12. 잠 15:33
13. 스 2:3
14. 잠 22:15
15. 잠 16:18
16. 잠 19:3, 잠 24:9
17. 잠 22:17-21
18. 전 12:13
19. 시 119:38, 시 119:4
"""

# 성경의 책 이름 리스트
bible_books = [
    "창세기", "출애굽기", "레위기", "민수기", "신명기", "여호수아", "사사기", "룻기", "사무엘상", "사무엘하",
    "열왕기상", "열왕기하", "역대상", "역대하", "에스라", "느헤미야", "에스더", "욥기", "시편", "잠언",
    "전도서", "아가", "이사야", "예레미야", "예레미야애가", "에스겔", "다니엘", "호세아", "요엘", "아모스",
    "오바댜", "요나", "미가", "나훔", "하박국", "스바냐", "학개", "스가랴", "말라기", "마태복음", "마가복음",
    "누가복음", "요한복음", "사도행전", "로마서", "고린도전서", "고린도후서", "갈라디아서", "에베소서", "빌립보서", "골로새서",
    "데살로니가전서", "데살로니가후서", "디모데전서", "디모데후서", "디도서", "빌레몬서", "히브리서", "야고보서", "베드로전서", "베드로후서",
    "요한일서", "요한이서", "요한삼서", "유다서", "요한계시록"
]

book_abbr_map = {'창': '창세기','출': '출애굽기','레': '레위기','민': '민수기','신': '신명기','수': '여호수아','삿': '사사기','룻': '룻기','삼상': '사무엘상','삼하': '사무엘하','왕상': '열왕기상','왕하': '열왕기하','대상': '역대상','대하': '역대하','스': '에스라','느': '느헤미야','에': '에스더','욥': '욥기','시': '시편','잠': '잠언','전': '전도서','아': '아가','사': '이사야','렘': '예레미야','애': '예레미야애가','겔': '에스겔','단': '다니엘','호': '호세아','욜': '요엘','암': '아모스','옵': '오바댜','욘': '요나','미': '미가','나': '나훔','합': '하박국','습': '스바냐','학': '학개','슥': '스가랴','말': '말라기','마': '마태복음','막': '마가복음','눅': '누가복음','요': '요한복음','행': '사도행전','롬': '로마서','고전': '고린도전서','고후': '고린도후서','갈': '갈라디아서','엡': '에베소서','빌': '빌립보서','골': '골로새서','살전': '데살로니가전서','살후': '데살로니가후서','딤전': '디모데전서','딤후': '디모데후서','딛': '디도서','몬': '빌레몬서','히': '히브리서','약': '야고보서','벧전': '베드로전서','벧후': '베드로후서','요일': '요한일서','요이': '요한이서','요삼': '요한삼서','유': '유다서','계': '요한계시록'
}

book_abbr_map_eng = {
    "창": "Gen.", "출": "Exod.", "레": "Lev.", "민": "Num.", "신": "Deut.",
    "수": "Josh.", "삿": "Judg.", "룻": "Ruth.", "삼상": "1Sam.", "삼하": "2Sam.",
    "열상": "1Kings.", "열하": "2Kings.", "대상": "1Chr.", "대하": "2Chr.",
    "에스라": "Ezra.", "느헤미야": "Neh.", "에스더": "Esth.", "욥": "Job.", "시": "Ps.",
    "잠": "Prov.", "전": "Eccles.", "아가": "Song.", "사": "Isa.", "렘": "Jer.",
    "애": "Lam.", "겔": "Ezek.", "단": "Dan.", "호": "Hos.", "욜": "Joel.",
    "암": "Amos.", "옵": "Obad.", "욘": "Jonah.", "미": "Mic.", "나": "Nah.",
    "합": "Hab.", "스": "Zeph.", "학": "Hag.", "슥": "Zech.", "말": "Mal.",
    "마": "Matt.", "막": "Mark.", "눅": "Luke.", "요": "John.", "행": "Acts.",
    "롬": "Rom.", "고전": "1Cor.", "고후": "2Cor.", "갈": "Gal.", "엡": "Eph.",
    "빌": "Phil.", "골": "Col.", "살전": "1Thess.", "살후": "2Thess.", "딤전": "1Tim.",
    "딤후": "2Tim.", "Tit": "Titus.", "몬": "Phlm.", "히": "Heb.", "약": "Jas.",
    "벧전": "1Pet.", "벧후": "2Pet.", "요일": "1John", "요이": "2John", "요삼": "3John",
    "유": "Jude.", "계": "Rev."
}

# '개역개정-txt' 폴더 내 파일을 읽어옴
def read_files_in_directory(directory):
    file_contents = []
    for filename in os.listdir(directory):
        if filename.endswith('.txt'):
            with open(os.path.join(directory, filename), 'r', encoding='utf-8') as file:
                content = file.read()
                file_contents.append(content)
    return file_contents

texts = read_files_in_directory('개역개정-text')

# texts 리스트를 딕셔너리로 변환
bible_dict = {bible_books[i]: texts[i] for i in range(len(bible_books))}


def split_and_format_verses(bible_dict):
    result = {}

    for book, verses in bible_dict.items():
        chapter_map = defaultdict(list)

        for verse in verses.splitlines():
            match = re.match(r'([가-힣]+)(\d+):(\d+)\s+(.*)', verse)
            if match:
                chapter = match.group(2)  # 장 번호
                verse_num = match.group(3)  # 절 번호
                content = match.group(4)  # 본문
                chapter_map[chapter].append(f"{verse_num} {content}")

        # 장 번호 순서대로 정렬 (숫자 기준)
        sorted_chapters = sorted(chapter_map.items(), key=lambda x: int(x[0]))
        chapter_list = [verses for _, verses in sorted_chapters]
        result[book] = chapter_list

    return result

formatted_bible = split_and_format_verses(bible_dict)

def parse_multi_refs_line(text):
    """번호 + 여러 구절이 한 줄에 포함된 경우 처리"""
    lines = text.strip().split('\n')
    grouped_refs = []

    for line in lines:
        parts = line.strip().split(' ', 1)
        if len(parts) < 2:
            continue

        ref_text = parts[1]
        ref_items = [r.strip() for r in ref_text.split(',')]
        grouped_refs.append(ref_items)

    return grouped_refs

grouped_refs = parse_multi_refs_line(raw_text)

def extract_passages_grouped(data, grouped_refs):
    result = []

    for ref_group in grouped_refs:
        merged_verses = []
        merged_label = []

        for ref in ref_group:
            match = re.match(r'([가-힣]+)\s+(\d+):([\d\-]+)', ref)
            if not match:
                continue
            abbr, chapter, verses = match.groups()
            book = book_abbr_map.get(abbr, abbr)
            chapter_idx = int(chapter) - 1
            chapter_data = data.get(book, [])

            if chapter_idx >= len(chapter_data):
                continue

            chapter_content = chapter_data[chapter_idx]

            if '-' in verses:
                start, end = map(int, verses.split('-'))
                if end > len(chapter_content):
                    continue
                verse_text = '\n'.join(chapter_content[v - 1] for v in range(start, end + 1))
                merged_label.append(f"{book} {chapter}:{start}-{end}\n")
                merged_verses.append(verse_text)
            else:
                v = int(verses)
                if v > len(chapter_content):
                    continue
                verse_text = chapter_content[v - 1]
                merged_label.append(f"{book} {chapter}:{v}\n")
                merged_verses.append(verse_text+"\n")

        # 최종 병합
        label = ''.join(merged_label)
        content = ''.join(merged_verses)
        result.append([label, content])

    return result

extracted = extract_passages_grouped(formatted_bible, grouped_refs)

def parse_scripture_file(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()

    result = defaultdict(lambda: defaultdict(list))
    
    # 책 이름 뒤에 공백이 있는 경우를 반영한 정규표현식
    pattern = re.compile(r'^([A-Za-z0-9]+\.?)\s+(\d+):(\d+)\s+(.*)')

    for line in lines:
        match = pattern.match(line.strip())
        if match:
            book, chapter, verse, content = match.groups()
            chapter = int(chapter)
            result[book][chapter].append(f"{verse} {content.strip()}")

    final_result = {}
    for book, chapters in result.items():
        max_chapter = max(chapters)
        chapter_list = [chapters[i] if i in chapters else [] for i in range(1, max_chapter + 1)]
        final_result[book] = chapter_list

    return final_result

def extract_passages_grouped_eng(data, grouped_refs):
    result = []

    for ref_group in grouped_refs:
        merged_verses = []
        merged_label = []

        for ref in ref_group:
            match = re.match(r'([가-힣]+)\s+(\d+):([\d\-]+)', ref)
            if not match:
                continue
            abbr, chapter, verses = match.groups()
            book = book_abbr_map_eng.get(abbr, abbr)
            chapter_idx = int(chapter) - 1
            chapter_data = data.get(book, [])

            if chapter_idx >= len(chapter_data):
                print("비상!!!!!!!!!!!")
                continue

            chapter_content = chapter_data[chapter_idx]

            if '-' in verses:
                start, end = map(int, verses.split('-'))
                if end > len(chapter_content):
                    continue
                verse_text = '\n'.join(chapter_content[v - 1] for v in range(start, end + 1))
                merged_label.append(f"{abbr} {chapter}:{start}-{end}\n")
                merged_verses.append(verse_text)
            else:
                v = int(verses)
                if v > len(chapter_content):
                    continue
                verse_text = chapter_content[v - 1]
                merged_label.append(f"{book} {chapter}:{v}\n")
                merged_verses.append(verse_text+"\n")

        # 최종 병합
        label = ''.join(merged_label)
        content = ''.join(merged_verses)
        result.append([label, content])

    return result

parsed = parse_scripture_file("KJV-text/KJV_text.txt")

extracted_passages_eng = extract_passages_grouped_eng(parsed, grouped_refs)

from pptx import Presentation
from pptx.util import Pt
import copy
import os
from pptx.dml.color import RGBColor

def add_scripture_to_ppt(template_path, verse_texts, output_path="output.pptx"):
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {template_path}")
    
    prs = Presentation(template_path)

    # 빈 레이아웃 (일반적으로 '빈 화면'은 레이아웃 인덱스 6번입니다)
    blank_layout = prs.slide_layouts[6]

    def duplicate_slide_with_blank_layout(prs, slide):
        new_slide = prs.slides.add_slide(blank_layout)  # 빈 레이아웃 사용
        for shape in slide.shapes:
            new_shape = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')
        return new_slide

    # 슬라이드 수 부족하면 복제
    while len(prs.slides) < len(verse_texts):
        duplicate_slide_with_blank_layout(prs, prs.slides[-1])

    for idx, (address, verse) in enumerate(verse_texts):
        slide = prs.slides[idx]

        # 2번째 텍스트 상자 (인덱스 1)에 본문 텍스트 추가
        text_shape = slide.shapes[1]
        text_frame = text_shape.text_frame
        text_frame.clear()
        for i, line in enumerate(verse.split('\n')):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = line
            p.font.color.rgb = RGBColor(31, 51, 55)
            p.font.size = Pt(28)
            p.font.name = '나눔스퀘어 네오 Bold'

        # 3번째 텍스트 상자 (인덱스 2)에 주소 텍스트 추가
        text_shape = slide.shapes[2]
        text_frame = text_shape.text_frame
        text_frame.clear()
        for i, line in enumerate(address.split('\n')):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = line
            p.font.color.rgb = RGBColor(31, 51, 55)
            p.font.size = Pt(37.3)
            p.font.name = '나눔스퀘어 네오 ExtraBold'

        prs.save(output_path)

    for idx, (address, verse) in enumerate(extracted_passages_eng):
        slide = prs.slides[idx]

        # 2번째 텍스트 상자 (인덱스 1)에 본문 텍스트 추가
        text_shape = slide.shapes[6]
        text_frame = text_shape.text_frame
        text_frame.clear()
        for i, line in enumerate(address.split('\n')):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = line
            p.font.color.rgb = RGBColor(143, 167, 159)
            p.font.size = Pt(28)
            p.font.name = '나눔스퀘어 네오 ExtraBold'


        # 2번째 텍스트 상자 (인덱스 1)에 본문 텍스트 추가
        text_shape = slide.shapes[7]
        text_frame = text_shape.text_frame
        text_frame.clear()
        for i, line in enumerate(verse.split('\n')):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = line
            p.font.color.rgb = RGBColor(79, 101, 94)
            p.font.size = Pt(28)
            p.font.name = 'Pretendard Variable'

        prs.save(output_path)

    print(f"저장 완료: {output_path}")

# 예시 사용법
template_path = "template.pptx"  # 템플릿 PPTX 파일 경로
output_path = "output.pptx"  # 출력할 PPTX 파일 경로
add_scripture_to_ppt(template_path, extracted, output_path)