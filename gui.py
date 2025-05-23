import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os
import re
import copy
from collections import defaultdict
import sys

"""
exe 생성 빌드
pyinstaller --noconfirm --onefile --add-data "개역개정-text;개역개정-text" --add-data "ESV-text/ESV_cleaned.txt;ESV-text" gui.py
"""

# 실행 경로 얻기 (PyInstaller 환경과 일반 환경 모두 지원)
def resource_path(relative_path):
    if hasattr(sys, '_MEIPASS'):
        # PyInstaller 임시폴더 경로
        base_path = sys._MEIPASS
    else:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)


# ------------------------- 성경 데이터 처리 -------------------------

def read_files_in_directory(directory):
    file_contents = []
    for filename in os.listdir(directory):
        if filename.endswith('.txt'):
            with open(os.path.join(directory, filename), 'r', encoding='utf-8') as file:
                content = file.read()
                file_contents.append(content)
    return file_contents

bible_books = [
    "창세기", "출애굽기", "레위기", "민수기", "신명기", "여호수아", "사사기", "룻기", "사무엘상", "사무엘하",
    "열왕기상", "열왕기하", "역대상", "역대하", "에스라", "느헤미야", "에스더", "욥기", "시편", "잠언",
    "전도서", "아가", "이사야", "예레미야", "예레미야애가", "에스겔", "다니엘", "호세아", "요엘", "아모스",
    "오바댜", "요나", "미가", "나훔", "하박국", "스바냐", "학개", "스가랴", "말라기", "마태복음", "마가복음",
    "누가복음", "요한복음", "사도행전", "로마서", "고린도전서", "고린도후서", "갈라디아서", "에베소서", "빌립보서", "골로새서",
    "데살로니가전서", "데살로니가후서", "디모데전서", "디모데후서", "디도서", "빌레몬서", "히브리서", "야고보서", "베드로전서", "베드로후서",
    "요한일서", "요한이서", "요한삼서", "유다서", "요한계시록"
]

book_abbr_map = {'창': '창세기','출': '출애굽기','레': '레위기','민': '민수기','신': '신명기','수': '여호수아','삿': '사사기','룻': '룻기','삼상': '사무엘상','삼하': '사무엘하','왕상': '열왕기상','왕하': '열왕기하','대상': '역대상','대하': '역대하','스': '에스라','느': '느헤미야','에': '에스더','욥': '욥기','시': '시편','잠': '잠언','전': '전도서','아': '아가','사': '이사야','렘': '예레미야','애': '예레미야애가','겔': '에스겔','단': '다니엘','호': '호세아','욜': '요엘','암': '아모스','옵': '오바댜','욘': '요나','미': '미가','나': '나훔','합': '하박국','습': '스바냐','학': '학개','슥': '스가랴','말': '말라기','마': '마태복음','막': '마가복음','눅': '누가복음','요': '요한복음','행': '사도행전','롬': '로마서','고전': '고린도전서','고후': '고린도후서','갈': '갈라디아서','엡': '에베소서','빌': '빌립보서','골': '골로새서','살전': '데살로니가전서','살후': '데살로니가후서','딤전': '디모데전서','딤후': '디모데후서','딛': '디도서','몬': '빌레몬서','히': '히브리서','약': '야고보서','벧전': '베드로전서','벧후': '베드로후서','요일': '요한일서','요이': '요한이서','요삼': '요한삼서','유': '유다서','계': '요한계시록'}

# KJV 성경 매핑
book_abbr_map_eng = {
    "창": "Gen.", "출": "Exod.", "레": "Lev.", "민": "Num.", "신": "Deut.",
    "수": "Josh.", "삿": "Judg.", "룻": "Ruth.", "삼상": "1Sam.", "삼하": "2Sam.",
    "열상": "1Kings.", "열하": "2Kings.", "대상": "1Chr.", "대하": "2Chr.",
    "에스라": "Ezra.", "느헤미야": "Neh.", "에스더": "Esth.", "욥": "Job.", "시": "Ps.",
    "잠": "Prov.", "전": "Eccles.", "아가": "Song.", "사": "Isa.", "렘": "Jer.",
    "애": "Lam.", "겔": "Ezek.", "단": "Dan.", "호": "Hos.", "욜": "Joel.",
    "암": "Amos.", "옵": "Obad.", "욘": "Jonah.", "미": "Mic.", "나": "Nah.",
    "합": "Hab.", "스": "Zeph.", "학": "Hag.", "슥": "Zech.", "말": "Mal.",
    "마": "Matt.", "막": "Mark.", "눅": "Luke.", "요": "John.", "행": "Acts.",
    "롬": "Rom.", "고전": "1Cor.", "고후": "2Cor.", "갈": "Gal.", "엡": "Eph.",
    "빌": "Phil.", "골": "Col.", "살전": "1Thess.", "살후": "2Thess.", "딤전": "1Tim.",
    "딤후": "2Tim.", "Tit": "Titus.", "몬": "Phlm.", "히": "Heb.", "약": "Jas.",
    "벧전": "1Pet.", "벧후": "2Pet.", "요일": "1John", "요이": "2John", "요삼": "3John",
    "유": "Jude.", "계": "Rev."
}

# ESV 성경 매핑
bible_book_abbreviations = {
    '창': 'Gen', '출': 'Exo', '레': 'Lev', '민': 'Num', '신': 'Deu', '수': 'Jos', '삿': 'Jdg', '룻': 'Rth', '삼상': '1Sa', '삼하': '2Sa', '왕상': '1Ki', '왕하': '2Ki', '대상': '1Ch', '대하': '2Ch', '스': 'Ezr', '느': 'Neh', '에': 'Est', '욥': 'Job', '시': 'Psa', '잠': 'Pro', '전': 'Ecc', '아': 'Son', '사': 'Isa', '렘': 'Jer', '애': 'Lam', '겔': 'Eze', '단': 'Dan', '호': 'Hos', '욜': 'Joe', '암': 'Amo', '옵': 'Oba', '욘': 'Jon', '미': 'Mic', '나': 'Nah', '합': 'Hab', '습': 'Zep', '학': 'Hag', '슥': 'Zec', '말': 'Mal', '마': 'Mat', '막': 'Mar', '눅': 'Luk', '요': 'Joh', '행': 'Act', '롬': 'Rom', '고전': '1Co', '고후': '2Co', '갈': 'Gal', '엡': 'Eph', '빌': 'Php', '골': 'Col', '살전': '1Th', '살후': '2Th', '딤전': '1Ti', '딤후': '2Ti', '딛': 'Tit', '몬': 'Phm', '히': 'Heb', '약': 'Jam', '벧전': '1Pe', '벧후': '2Pe', '요일': '1Jo', '요이': '2Jo', '요삼': '3Jo', '유': 'Jud', '계': 'Rev'
}

# texts = read_files_in_directory('개역개정-text')
# bible_dict = {bible_books[i]: texts[i] for i in range(len(bible_books))}
texts = read_files_in_directory(resource_path('개역개정-text'))
bible_dict = {bible_books[i]: texts[i] for i in range(len(bible_books))}

def split_and_format_verses(bible_dict):
    result = {}
    for book, verses in bible_dict.items():
        chapter_map = defaultdict(list)
        for verse in verses.splitlines():
            match = re.match(r'([가-힣]+)(\d+):(\d+)\s+(.*)', verse)
            if match:
                chapter = match.group(2)
                verse_num = match.group(3)
                content = match.group(4)
                chapter_map[chapter].append(f"{verse_num} {content}")
        sorted_chapters = sorted(chapter_map.items(), key=lambda x: int(x[0]))
        chapter_list = [verses for _, verses in sorted_chapters]
        result[book] = chapter_list
    return result

formatted_bible = split_and_format_verses(bible_dict)

# ------------------------- GUI 구성 -------------------------

def parse_multi_refs_line(text):
    lines = text.strip().split('\n')
    grouped_refs = []
    for line in lines:
        parts = line.strip().split(' ', 1)
        if len(parts) < 2:
            continue
        ref_text = parts[1]
        ref_items = [r.strip() for r in ref_text.split(';')]
        grouped_refs.append(ref_items)
    return grouped_refs

def extract_passages_grouped(data, grouped_refs):
    result = []
    for ref_group in grouped_refs:
        merged_verses = []
        merged_label = []
        
        for ref in ref_group:
            if ref[:5] == '<인용구>':
                merged_label.append(f"{ref[:5]}\n")
                merged_verses.append(ref[5:])
                continue
        for ref in ref_group:
            match = re.match(r'([가-힣]+)\s+(\d+):([\d\-]+)', ref)
            if not match:
                continue
            abbr, chapter, verses = match.groups()
            book = book_abbr_map.get(abbr, abbr)
            chapter_idx = int(chapter) - 1
            chapter_data = data.get(book, [])
            if chapter_idx >= len(chapter_data):
                continue
            chapter_content = chapter_data[chapter_idx]
            if '-' in verses:
                start, end = map(int, verses.split('-'))
                if end > len(chapter_content):
                    continue
                verse_text = ' '.join(chapter_content[v - 1] for v in range(start, end + 1))
                merged_label.append(f"{book} {chapter}:{start}-{end}\n")
                merged_verses.append(verse_text)
            else:
                v = int(verses)
                if v > len(chapter_content):
                    continue
                verse_text = chapter_content[v - 1]
                merged_label.append(f"{book} {chapter}:{v}\n")
                merged_verses.append(verse_text)
        label = ''.join(merged_label)
        content = ' '.join(merged_verses)
        result.append((label, content))
    return result

def parse_scripture_file(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()

    result = defaultdict(lambda: defaultdict(list))
    
    # 책 이름 뒤에 공백이 있는 경우를 반영한 정규표현식
    pattern = re.compile(r'^([A-Za-z0-9]+\.?)\s+(\d+):(\d+)\s+(.*)')

    for line in lines:
        match = pattern.match(line.strip())
        if match:
            book, chapter, verse, content = match.groups()
            chapter = int(chapter)
            result[book][chapter].append(f"{verse} {content.strip()}")

    final_result = {}
    for book, chapters in result.items():
        max_chapter = max(chapters)
        chapter_list = [chapters[i] if i in chapters else [] for i in range(1, max_chapter + 1)]
        final_result[book] = chapter_list

    return final_result

def extract_passages_grouped_eng(data, grouped_refs):
    result = []

    for ref_group in grouped_refs:
        merged_verses = []
        merged_label = []

        for ref in ref_group:
            match = re.match(r'([가-힣]+)\s+(\d+):([\d\-]+)', ref)
            if not match:
                continue
            abbr, chapter, verses = match.groups()
            book = bible_book_abbreviations.get(abbr, abbr)
            chapter_idx = int(chapter) - 1
            chapter_data = data.get(book, [])

            if chapter_idx >= len(chapter_data):
                print("비상!!!!!!!!!!!")
                continue

            chapter_content = chapter_data[chapter_idx]

            if '-' in verses:
                start, end = map(int, verses.split('-'))
                if end > len(chapter_content):
                    continue
                verse_text = ' '.join(chapter_content[v - 1] for v in range(start, end + 1))
                merged_label.append(f"{book} {chapter}:{start}-{end}\n")
                merged_verses.append(verse_text)
            else:
                v = int(verses)
                if v > len(chapter_content):
                    continue
                verse_text = chapter_content[v - 1]
                merged_label.append(f"{book} {chapter}:{v}\n")
                merged_verses.append(verse_text)

        # 최종 병합
        label = ''.join(merged_label)
        content = ' '.join(merged_verses)
        result.append([label, content])

    return result

# parsed = parse_scripture_file(resource_path("KJV-text/KJV_text.txt"))
parsed = parse_scripture_file(resource_path("ESV-text/ESV_cleaned.txt"))

def add_scripture_to_ppt(template_path, verse_texts, verse_texts_eng, output_path="output.pptx"):
    if not os.path.exists(resource_path(template_path)):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {template_path}")
    
    prs = Presentation(template_path)

    # 빈 레이아웃 (일반적으로 '빈 화면'은 레이아웃 인덱스 6번입니다)
    blank_layout = prs.slide_layouts[6]

    def duplicate_slide_with_blank_layout(prs, slide):
        new_slide = prs.slides.add_slide(blank_layout)  # 빈 레이아웃 사용
        for shape in slide.shapes:
            new_shape = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')
        return new_slide

    # 슬라이드 수 부족하면 복제
    while len(prs.slides) < len(verse_texts):
        duplicate_slide_with_blank_layout(prs, prs.slides[-1])

    for idx, (address, verse) in enumerate(verse_texts):
        slide = prs.slides[idx]

        # 2번째 텍스트 상자 (인덱스 1)에 본문 텍스트 추가
        text_shape = slide.shapes[1]
        text_frame = text_shape.text_frame
        text_frame.clear()
        text_frame.text = verse  # 바로 텍스트를 할당하여 빈 문단 없이 설정
        text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 51, 55)  # 검정색으로 설정
        text_frame.paragraphs[0].font.size = Pt(28)  # 첫 문단의 폰트 설정
        text_frame.paragraphs[0].font.name = '나눔스퀘어 네오 Bold'

        # 3번째 텍스트 상자 (인덱스 2)에 주소 텍스트 추가
        text_shape = slide.shapes[2]
        text_frame = text_shape.text_frame
        text_frame.clear()
        for i, line in enumerate(address.split('\n')):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = line
            p.font.color.rgb = RGBColor(31, 51, 55)
            p.font.size = Pt(37.3)
            p.font.name = '나눔스퀘어 네오 ExtraBold'

        prs.save(output_path)

    for idx, (address, verse) in enumerate(verse_texts_eng):
        slide = prs.slides[idx]

        # 2번째 텍스트 상자 (인덱스 1)에 본문 텍스트 추가
        text_shape = slide.shapes[6]
        text_frame = text_shape.text_frame
        text_frame.clear()
        for i, line in enumerate(address.split('\n')):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = line
            p.font.color.rgb = RGBColor(143, 167, 159)
            p.font.size = Pt(28)
            p.font.name = '나눔스퀘어 네오 ExtraBold'


        # 2번째 텍스트 상자 (인덱스 1)에 본문 텍스트 추가
        text_shape = slide.shapes[7]
        text_frame = text_shape.text_frame
        text_frame.clear()
        text_frame.text = verse
        text_frame.paragraphs[0].font.size = Pt(20)
        text_frame.paragraphs[0].font.name = 'Pretendard Variable'
        text_frame.paragraphs[0].font.color.rgb = RGBColor(79, 101, 94)  # 검정색으로 설정

        prs.save(output_path)

def on_generate_click():
    raw_text = input_text.get("1.0", tk.END)
    grouped_refs = parse_multi_refs_line(raw_text)
    extracted = extract_passages_grouped(formatted_bible, grouped_refs)
    extracted_eng = extract_passages_grouped_eng(parsed, grouped_refs)
    if not extracted:
        messagebox.showerror("오류", "유효한 구절을 입력하세요.")
        return
    save_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
    # PPTX 파일 저장 후 자동 실행
    if save_path:
        add_scripture_to_ppt(template_path='template.pptx', verse_texts=extracted, verse_texts_eng=extracted_eng, output_path=save_path)
        os.startfile(save_path)
        # messagebox.showinfo("완료", f"PPTX 파일이 저장되었습니다: {save_path}")

# ------------------------- 실행 -------------------------

root = tk.Tk()
root.title("성경 구절 PPTX 변환기")
tk.Label(root, text=(
    "구절 입력 양식\n"
    "1) 각 구절은 '숫자. 책 장:절' 형식으로 입력합니다.\n"
    "2) 한 슬라이드에 여러 구절을 넣고 싶다면, 각 구절을 세미콜론(;)으로 구분합니다.\n"
    "3) 이어진 구절로 이루어져 있다면 '책 장:절-절' 형식으로 입력합니다.\n"
    "4) 책 제목은 약어로 입력합니다.\n\n"
    "5) 인용구는 '<인용구> 내용' 형식으로 입력합니다.\n\n"
    "예)\n"
    "1. 잠 1:8; 잠 31:2\n"
    "2. 마 5:7-9\n"
    "3. <인용구> \"동해물과 백두산이 마르고 닳도록\"\n"
)).pack(pady=5)
input_text = tk.Text(root, height=15, width=60)
input_text.pack(padx=10)
tk.Button(root, text="PPTX로 변환", command=on_generate_click).pack(pady=10)
root.mainloop()